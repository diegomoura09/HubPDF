"""
Comprehensive PDF conversion service with Office, image, and text conversion support
"""
import os
import tempfile
import subprocess
import logging
import asyncio
from pathlib import Path
from typing import List, Optional, Dict, Any, Tuple
import uuid
import time
from concurrent.futures import ThreadPoolExecutor
import shutil

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    import pikepdf
except ImportError:
    pikepdf = None

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from PIL import Image, ImageDraw, ImageFont
except ImportError:
    Image = ImageDraw = ImageFont = None

try:
    from pdfminer.high_level import extract_text as pdfminer_extract_text
    from pdfminer.layout import LAParams
except ImportError:
    pdfminer_extract_text = LAParams = None

try:
    import pytesseract
except ImportError:
    pytesseract = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from openpyxl import Workbook, load_workbook
except ImportError:
    Workbook = load_workbook = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.pagesizes import letter
except ImportError:
    SimpleDocTemplate = Paragraph = Spacer = getSampleStyleSheet = letter = None

from app.config import settings
from app.pdf_compress import compress_pdf as compress_pdf_advanced

logger = logging.getLogger(__name__)

class ConversionError(Exception):
    """Custom exception for conversion errors"""
    pass

class ConversionService:
    """
    Comprehensive PDF conversion service supporting:
    - PDF ↔ Office formats (DOCX/XLSX/PPTX)
    - PDF ↔ Images (PNG/JPG/JPEG/ICO)
    - PDF → Text extraction with OCR fallback
    - PDF operations (split, compress, merge)
    """
    
    def __init__(self):
        self.temp_dir = Path(tempfile.gettempdir()) / "hubpdf_conversions"
        self.temp_dir.mkdir(exist_ok=True)
        self.executor = ThreadPoolExecutor(max_workers=4)
        
    def _get_work_dir(self, job_id: str) -> Path:
        """Create and return a unique working directory for the job"""
        work_dir = self.temp_dir / job_id
        work_dir.mkdir(parents=True, exist_ok=True)
        return work_dir
        
    def _cleanup_work_dir(self, work_dir: Path):
        """Clean up the working directory"""
        try:
            if work_dir.exists():
                shutil.rmtree(work_dir)
        except Exception as e:
            logger.warning(f"Failed to cleanup work directory {work_dir}: {e}")
    
    def _get_original_filename_with_suffix(self, file_path: str, suffix: str, new_ext: str) -> str:
        """Get filename preserving original name with suffix"""
        file_path = Path(file_path)
        original_stem = file_path.stem
        return f"{original_stem}_{suffix}.{new_ext}"
    
    def _run_libreoffice_conversion(self, input_path: str, output_dir: str, target_format: str, timeout: int = 60) -> str:
        """Run LibreOffice headless conversion with timeout"""
        try:
            cmd = [
                "soffice", 
                "--headless",
                "--invisible", 
                "--nodefault",
                "--norestore",
                "--nologo",
                f"--convert-to", target_format,
                f"--outdir", output_dir,
                input_path
            ]
            
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=timeout,
                cwd=output_dir
            )
            
            if result.returncode != 0:
                raise ConversionError(f"LibreOffice conversion failed: {result.stderr}")
                
            # Find the output file
            input_stem = Path(input_path).stem
            output_files = list(Path(output_dir).glob(f"{input_stem}.*"))
            
            if not output_files:
                raise ConversionError("LibreOffice conversion produced no output file")
                
            return str(output_files[0])
            
        except subprocess.TimeoutExpired:
            # Kill any hanging soffice processes
            subprocess.run(["pkill", "-f", "soffice"], capture_output=True)
            raise ConversionError("LibreOffice conversion timed out")
        except FileNotFoundError:
            raise ConversionError("LibreOffice not found. Please install libreoffice package.")
    
    # Office format conversions
    async def pdf_to_docx(self, pdf_path: str, job_id: str = None) -> str:
        """Convert PDF to DOCX using Python libraries"""
        if Document is None:
            raise ConversionError("python-docx not available. Please install python-docx package.")
        if job_id is None:
            job_id = str(uuid.uuid4())
            
        work_dir = self._get_work_dir(job_id)
        
        try:
            def _extract_and_create_docx():
                import pdfplumber
                
                # Create new DOCX document
                doc = Document()
                
                with pdfplumber.open(pdf_path) as pdf:
                    for page_num, page in enumerate(pdf.pages):
                        text = page.extract_text()
                        if text:
                            # Add page heading
                            if page_num > 0:
                                doc.add_page_break()
                            doc.add_heading(f'Page {page_num + 1}', level=2)
                            # Add text content
                            doc.add_paragraph(text)
                
                # Save with original filename + suffix
                output_filename = self._get_original_filename_with_suffix(pdf_path, "to_docx", "docx")
                output_path = work_dir / output_filename
                doc.save(str(output_path))
                return str(output_path)
            
            loop = asyncio.get_event_loop()
            output_path = await loop.run_in_executor(self.executor, _extract_and_create_docx)
            return output_path
            
        except Exception as e:
            self._cleanup_work_dir(work_dir)
            raise ConversionError(f"PDF to DOCX conversion failed: {str(e)}")
    
    async def docx_to_pdf(self, docx_path: str, job_id: str = None) -> str:
        """Convert DOCX to PDF using Python libraries"""
        if Document is None or SimpleDocTemplate is None:
            raise ConversionError("Required libraries not available. Please install python-docx and reportlab.")
        if job_id is None:
            job_id = str(uuid.uuid4())
            
        work_dir = self._get_work_dir(job_id)
        
        try:
            def _convert_docx_to_pdf():
                # Read DOCX content
                doc = Document(docx_path)
                
                # Create PDF with reportlab
                output_filename = self._get_original_filename_with_suffix(docx_path, "to_pdf", "pdf")
                output_path = work_dir / output_filename
                
                pdf_doc = SimpleDocTemplate(str(output_path), pagesize=letter)
                styles = getSampleStyleSheet()
                story = []
                
                # Extract text from DOCX paragraphs
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        p = Paragraph(paragraph.text, styles['Normal'])
                        story.append(p)
                        story.append(Spacer(1, 12))
                
                if not story:
                    # Add a placeholder if no content
                    story.append(Paragraph("Document converted from DOCX", styles['Normal']))
                
                pdf_doc.build(story)
                return str(output_path)
            
            loop = asyncio.get_event_loop()
            output_path = await loop.run_in_executor(self.executor, _convert_docx_to_pdf)
            return output_path
            
        except Exception as e:
            self._cleanup_work_dir(work_dir)
            raise ConversionError(f"DOCX to PDF conversion failed: {str(e)}")
    
    async def pdf_to_xlsx(self, pdf_path: str, job_id: str = None) -> str:
        """Convert PDF to XLSX using Python libraries"""
        if Workbook is None:
            raise ConversionError("openpyxl not available. Please install openpyxl package.")
        if job_id is None:
            job_id = str(uuid.uuid4())
            
        work_dir = self._get_work_dir(job_id)
        
        try:
            def _extract_and_create_xlsx():
                import pdfplumber
                
                # Create new Excel workbook
                wb = Workbook()
                
                with pdfplumber.open(pdf_path) as pdf:
                    for page_num, page in enumerate(pdf.pages):
                        # Create sheet for each page
                        if page_num == 0:
                            ws = wb.active
                            ws.title = f"Page {page_num + 1}"
                        else:
                            ws = wb.create_sheet(title=f"Page {page_num + 1}")
                        
                        # Extract tables if available
                        tables = page.extract_tables()
                        if tables:
                            row_offset = 1
                            for table in tables:
                                for row in table:
                                    for col_idx, cell in enumerate(row):
                                        ws.cell(row=row_offset, column=col_idx + 1, value=cell)
                                    row_offset += 1
                                row_offset += 1  # Add space between tables
                        else:
                            # Extract text if no tables
                            text = page.extract_text()
                            if text:
                                lines = text.split('\n')
                                for row_idx, line in enumerate(lines, start=1):
                                    ws.cell(row=row_idx, column=1, value=line)
                
                # Save with original filename + suffix
                output_filename = self._get_original_filename_with_suffix(pdf_path, "to_xlsx", "xlsx")
                output_path = work_dir / output_filename
                wb.save(str(output_path))
                return str(output_path)
            
            loop = asyncio.get_event_loop()
            output_path = await loop.run_in_executor(self.executor, _extract_and_create_xlsx)
            return output_path
            
        except Exception as e:
            self._cleanup_work_dir(work_dir)
            raise ConversionError(f"PDF to XLSX conversion failed: {str(e)}")
    
    async def xlsx_to_pdf(self, xlsx_path: str, job_id: str = None) -> str:
        """Convert XLSX to PDF using Python libraries"""
        if load_workbook is None or SimpleDocTemplate is None:
            raise ConversionError("Required libraries not available. Please install openpyxl and reportlab.")
        if job_id is None:
            job_id = str(uuid.uuid4())
            
        work_dir = self._get_work_dir(job_id)
        
        try:
            def _convert_xlsx_to_pdf():
                from reportlab.lib import colors
                from reportlab.lib.pagesizes import letter, landscape
                from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, PageBreak
                from reportlab.lib.styles import getSampleStyleSheet
                
                # Read Excel content
                wb = load_workbook(xlsx_path, data_only=True)
                
                # Create PDF with reportlab
                output_filename = self._get_original_filename_with_suffix(xlsx_path, "to_pdf", "pdf")
                output_path = work_dir / output_filename
                
                pdf_doc = SimpleDocTemplate(str(output_path), pagesize=landscape(letter))
                styles = getSampleStyleSheet()
                story = []
                
                # Process each sheet
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    
                    # Add sheet title
                    story.append(Paragraph(f"<b>{sheet_name}</b>", styles['Heading1']))
                    
                    # Extract data from sheet
                    data = []
                    for row in ws.iter_rows(values_only=True):
                        data.append([str(cell) if cell is not None else "" for cell in row])
                    
                    if data:
                        # Create table
                        table = Table(data)
                        table.setStyle(TableStyle([
                            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                            ('FONTSIZE', (0, 0), (-1, 0), 10),
                            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                            ('GRID', (0, 0), (-1, -1), 1, colors.black)
                        ]))
                        story.append(table)
                    
                    story.append(PageBreak())
                
                if not story:
                    story.append(Paragraph("Empty Excel document", styles['Normal']))
                
                pdf_doc.build(story)
                return str(output_path)
            
            loop = asyncio.get_event_loop()
            output_path = await loop.run_in_executor(self.executor, _convert_xlsx_to_pdf)
            return output_path
            
        except Exception as e:
            self._cleanup_work_dir(work_dir)
            raise ConversionError(f"XLSX to PDF conversion failed: {str(e)}")
    
    async def pdf_to_pptx(self, pdf_path: str, job_id: str = None) -> str:
        """Convert PDF to PPTX using Python libraries"""
        if Presentation is None or fitz is None:
            raise ConversionError("Required libraries not available. Please install python-pptx and pymupdf.")
        if job_id is None:
            job_id = str(uuid.uuid4())
            
        work_dir = self._get_work_dir(job_id)
        
        try:
            def _convert_pdf_to_pptx():
                from pptx.util import Inches
                
                # Create new PowerPoint presentation
                prs = Presentation()
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(7.5)
                
                # Convert each PDF page to image and add to presentation
                doc = fitz.open(pdf_path)
                
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    
                    # Convert page to image
                    mat = fitz.Matrix(2, 2)  # 2x zoom for better quality
                    pix = page.get_pixmap(matrix=mat, alpha=False)
                    
                    # Save temp image
                    img_path = work_dir / f"temp_page_{page_num}.png"
                    pix.save(str(img_path))
                    
                    # Add blank slide
                    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
                    slide = prs.slides.add_slide(blank_slide_layout)
                    
                    # Add image to slide
                    left = Inches(0.5)
                    top = Inches(0.5)
                    height = Inches(6.5)
                    slide.shapes.add_picture(str(img_path), left, top, height=height)
                    
                    # Clean up temp image
                    img_path.unlink(missing_ok=True)
                
                doc.close()
                
                # Save with original filename + suffix
                output_filename = self._get_original_filename_with_suffix(pdf_path, "to_pptx", "pptx")
                output_path = work_dir / output_filename
                prs.save(str(output_path))
                return str(output_path)
            
            loop = asyncio.get_event_loop()
            output_path = await loop.run_in_executor(self.executor, _convert_pdf_to_pptx)
            return output_path
            
        except Exception as e:
            self._cleanup_work_dir(work_dir)
            raise ConversionError(f"PDF to PPTX conversion failed: {str(e)}")
    
    async def pptx_to_pdf(self, pptx_path: str, job_id: str = None) -> str:
        """Convert PPTX to PDF using Python libraries"""
        if Presentation is None or Image is None or SimpleDocTemplate is None:
            raise ConversionError("Required libraries not available. Please install python-pptx, pillow, and reportlab.")
        if job_id is None:
            job_id = str(uuid.uuid4())
            
        work_dir = self._get_work_dir(job_id)
        
        try:
            def _convert_pptx_to_pdf():
                from reportlab.lib.pagesizes import letter, landscape
                from reportlab.pdfgen import canvas
                from reportlab.lib.utils import ImageReader
                
                # Read PowerPoint presentation
                prs = Presentation(pptx_path)
                
                # Create PDF with reportlab
                output_filename = self._get_original_filename_with_suffix(pptx_path, "to_pdf", "pdf")
                output_path = work_dir / output_filename
                
                # Create PDF with landscape orientation
                c = canvas.Canvas(str(output_path), pagesize=landscape(letter))
                page_width, page_height = landscape(letter)
                
                # Extract text from each slide and add to PDF
                for slide_num, slide in enumerate(prs.slides):
                    # Extract all text from slide
                    text_content = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            if shape.text.strip():
                                text_content.append(shape.text)
                    
                    # Add slide number
                    c.setFont("Helvetica-Bold", 16)
                    c.drawString(50, page_height - 50, f"Slide {slide_num + 1}")
                    
                    # Add text content
                    c.setFont("Helvetica", 12)
                    y_position = page_height - 100
                    for text in text_content:
                        # Wrap text if too long
                        max_width = page_width - 100
                        lines = []
                        words = text.split()
                        current_line = ""
                        
                        for word in words:
                            test_line = current_line + " " + word if current_line else word
                            if c.stringWidth(test_line, "Helvetica", 12) < max_width:
                                current_line = test_line
                            else:
                                if current_line:
                                    lines.append(current_line)
                                current_line = word
                        if current_line:
                            lines.append(current_line)
                        
                        for line in lines:
                            if y_position < 50:
                                c.showPage()
                                c.setFont("Helvetica", 12)
                                y_position = page_height - 50
                            c.drawString(50, y_position, line)
                            y_position -= 20
                    
                    c.showPage()
                
                c.save()
                return str(output_path)
            
            loop = asyncio.get_event_loop()
            output_path = await loop.run_in_executor(self.executor, _convert_pptx_to_pdf)
            return output_path
            
        except Exception as e:
            self._cleanup_work_dir(work_dir)
            raise ConversionError(f"PPTX to PDF conversion failed: {str(e)}")
    
    # Image conversions
    async def pdf_to_images(self, pdf_path: str, fmt: str = "png", dpi: int = 200, job_id: str = None) -> List[str]:
        """Convert PDF pages to images using PyMuPDF"""
        if fitz is None:
            raise ConversionError("PyMuPDF not available. Please install pymupdf package.")
            
        if job_id is None:
            job_id = str(uuid.uuid4())
            
        work_dir = self._get_work_dir(job_id)
        fmt = fmt.lower()
        
        if fmt not in ['png', 'jpg', 'jpeg']:
            raise ConversionError(f"Unsupported image format: {fmt}")
        
        try:
            def _convert_pages():
                doc = fitz.open(pdf_path)
                image_paths = []
                
                # Get original filename for prefix
                original_name = Path(pdf_path).stem
                
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    
                    # Create matrix for DPI scaling
                    mat = fitz.Matrix(dpi / 72, dpi / 72)
                    pix = page.get_pixmap(matrix=mat, alpha=False)
                    
                    # Save image with original filename prefix
                    if fmt in ['jpg', 'jpeg']:
                        image_path = work_dir / f"{original_name}_page_{page_num + 1:03d}.jpg"
                        pix.save(str(image_path), output="jpeg")
                    else:
                        image_path = work_dir / f"{original_name}_page_{page_num + 1:03d}.png"
                        pix.save(str(image_path), output="png")
                    
                    image_paths.append(str(image_path))
                    
                doc.close()
                return image_paths
            
            loop = asyncio.get_event_loop()
            image_paths = await loop.run_in_executor(self.executor, _convert_pages)
            return image_paths
            
        except Exception as e:
            self._cleanup_work_dir(work_dir)
            raise ConversionError(f"PDF to images conversion failed: {str(e)}")
    
    async def images_to_pdf(self, image_paths: List[str], job_id: str = None) -> str:
        """Convert multiple images to PDF"""
        if Image is None:
            raise ConversionError("Pillow not available. Please install pillow package.")
            
        if not image_paths:
            raise ConversionError("No images provided")
            
        if job_id is None:
            job_id = str(uuid.uuid4())
            
        work_dir = self._get_work_dir(job_id)
        # Create filename from first image
        first_image_name = Path(image_paths[0]).stem if image_paths else "images"
        output_filename = f"{first_image_name}_combined.pdf"
        output_path = work_dir / output_filename
        
        try:
            def _combine_images():
                images = []
                for img_path in image_paths:
                    img = Image.open(img_path)
                    # Convert to RGB if necessary
                    if img.mode != 'RGB':
                        img = img.convert('RGB')
                    images.append(img)
                
                if images:
                    # Save as PDF
                    images[0].save(
                        str(output_path),
                        "PDF",
                        resolution=200.0,
                        save_all=True,
                        append_images=images[1:] if len(images) > 1 else []
                    )
                    
                # Close images to free memory
                for img in images:
                    img.close()
                    
                return str(output_path)
            
            loop = asyncio.get_event_loop()
            result_path = await loop.run_in_executor(self.executor, _combine_images)
            return result_path
            
        except Exception as e:
            self._cleanup_work_dir(work_dir)
            raise ConversionError(f"Images to PDF conversion failed: {str(e)}")
    
    # Text extraction
    async def pdf_to_txt(self, pdf_path: str, use_ocr: bool = False, job_id: str = None) -> str:
        """Extract text from PDF with optional OCR fallback"""
        if job_id is None:
            job_id = str(uuid.uuid4())
            
        work_dir = self._get_work_dir(job_id)
        # Use original filename with suffix
        output_filename = self._get_original_filename_with_suffix(pdf_path, "extracted", "txt")
        output_path = work_dir / output_filename
        
        try:
            def _extract_text():
                text = ""
                
                # Try PyMuPDF first (fastest)
                if fitz:
                    try:
                        doc = fitz.open(pdf_path)
                        for page in doc:
                            page_text = page.get_text()
                            if page_text.strip():
                                text += page_text + "\n\n"
                        doc.close()
                    except Exception as e:
                        logger.warning(f"PyMuPDF extraction failed: {e}")
                
                # Fallback to pdfminer if no text found
                if not text.strip() and pdfminer_extract_text:
                    try:
                        text = pdfminer_extract_text(pdf_path, laparams=LAParams())
                    except Exception as e:
                        logger.warning(f"PDFMiner extraction failed: {e}")
                
                # OCR fallback if enabled and no text found
                if use_ocr and not text.strip() and pytesseract and fitz:
                    try:
                        logger.info("Attempting OCR extraction...")
                        doc = fitz.open(pdf_path)
                        for page_num in range(min(5, len(doc))):  # Limit OCR to first 5 pages
                            page = doc.load_page(page_num)
                            pix = page.get_pixmap()
                            img_data = pix.tobytes("png")
                            
                            # Save temporary image for OCR
                            temp_img_path = work_dir / f"temp_page_{page_num}.png"
                            with open(temp_img_path, "wb") as f:
                                f.write(img_data)
                            
                            # Run OCR
                            ocr_text = pytesseract.image_to_string(str(temp_img_path), lang='por+eng')
                            if ocr_text.strip():
                                text += f"[Page {page_num + 1}]\n{ocr_text}\n\n"
                            
                            # Cleanup temp image
                            temp_img_path.unlink(missing_ok=True)
                        doc.close()
                    except Exception as e:
                        logger.warning(f"OCR extraction failed: {e}")
                
                # Save extracted text
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write(text if text.strip() else "No text could be extracted from this PDF.")
                
                return str(output_path)
            
            loop = asyncio.get_event_loop()
            result_path = await loop.run_in_executor(self.executor, _extract_text)
            return result_path
            
        except Exception as e:
            self._cleanup_work_dir(work_dir)
            raise ConversionError(f"Text extraction failed: {str(e)}")
    
    # PDF operations
    async def merge_pdfs(self, pdf_paths: List[str], job_id: str = None) -> str:
        """Merge multiple PDFs into a single file preserving original names"""
        if not pdf_paths:
            raise ConversionError("No PDF files provided")
            
        if fitz is None:
            raise ConversionError("PyMuPDF not available. Please install pymupdf package.")
            
        if job_id is None:
            job_id = str(uuid.uuid4())
            
        work_dir = self._get_work_dir(job_id)
        
        try:
            def _merge_pdfs():
                # Create output filename from first PDF
                first_pdf_name = Path(pdf_paths[0]).stem if pdf_paths else "merged"
                output_filename = f"{first_pdf_name}_merge.pdf"
                output_path = work_dir / output_filename
                
                # Create merged document
                merged_doc = fitz.open()
                
                for pdf_path in pdf_paths:
                    try:
                        source_doc = fitz.open(pdf_path)
                        merged_doc.insert_pdf(source_doc)
                        source_doc.close()
                    except Exception as e:
                        logger.warning(f"Error merging {pdf_path}: {e}")
                        continue
                
                # Save merged PDF
                merged_doc.save(str(output_path))
                merged_doc.close()
                
                return str(output_path)
            
            loop = asyncio.get_event_loop()
            result_path = await loop.run_in_executor(self.executor, _merge_pdfs)
            return result_path
            
        except Exception as e:
            self._cleanup_work_dir(work_dir)
            raise ConversionError(f"PDF merge failed: {str(e)}")
    
    async def split_pdf(self, pdf_path: str, ranges: str, job_id: str = None) -> List[str]:
        """Split PDF into multiple files based on page ranges"""
        if fitz is None:
            raise ConversionError("PyMuPDF not available. Please install pymupdf package.")
            
        if job_id is None:
            job_id = str(uuid.uuid4())
            
        work_dir = self._get_work_dir(job_id)
        
        try:
            def _split_pdf():
                doc = fitz.open(pdf_path)
                total_pages = len(doc)
                output_files = []
                
                # Get original filename for prefix
                original_name = Path(pdf_path).stem
                
                # Parse ranges (e.g., "1-3,5,7-10")
                page_ranges = []
                for range_part in ranges.split(','):
                    range_part = range_part.strip()
                    if '-' in range_part:
                        start, end = map(int, range_part.split('-'))
                        page_ranges.append((start - 1, end - 1))  # Convert to 0-based
                    else:
                        page_num = int(range_part) - 1  # Convert to 0-based
                        page_ranges.append((page_num, page_num))
                
                # Create split documents with original filename
                for i, (start, end) in enumerate(page_ranges):
                    if start < 0 or end >= total_pages or start > end:
                        continue
                        
                    split_doc = fitz.open()
                    split_doc.insert_pdf(doc, from_page=start, to_page=end)
                    
                    # Use original filename with split suffix
                    output_filename = f"{original_name}_split_{i + 1}_pages_{start + 1}-{end + 1}.pdf"
                    output_path = work_dir / output_filename
                    split_doc.save(str(output_path))
                    split_doc.close()
                    
                    output_files.append(str(output_path))
                
                doc.close()
                return output_files
            
            loop = asyncio.get_event_loop()
            result_files = await loop.run_in_executor(self.executor, _split_pdf)
            return result_files
            
        except Exception as e:
            self._cleanup_work_dir(work_dir)
            raise ConversionError(f"PDF split failed: {str(e)}")
    
    async def compress_pdf(self, pdf_path: str, level: str = "balanced", job_id: str = None, 
                          grayscale: bool = False, rasterize: bool = False) -> Dict[str, Any]:
        """Compress PDF file using Ghostscript with advanced compression
        
        Compression levels:
        - light: Alta qualidade, ganho moderado (~20-30%)
        - balanced: Melhor custo/benefício (~40-60%)
        - strong: Redução máxima, pode perder qualidade (~60-80%)
        
        Options:
        - grayscale: Converte para tons de cinza
        - rasterize: Modo extremo (rasteriza páginas, perde seleção de texto)
        """
        if job_id is None:
            job_id = str(uuid.uuid4())
            
        work_dir = self._get_work_dir(job_id)
        output_filename = self._get_original_filename_with_suffix(pdf_path, "compressed", "pdf")
        output_path = work_dir / output_filename
        
        try:
            # Mapear níveis antigos para novos
            level_map = {
                "normal": "balanced",
                "medium": "balanced",
                "high": "strong",
                "maximum": "strong"
            }
            level = level_map.get(level, level)
            
            def _compress_pdf():
                return compress_pdf_advanced(
                    input_path=pdf_path,
                    output_path=str(output_path),
                    level=level,
                    grayscale=grayscale,
                    rasterize=rasterize
                )
            
            loop = asyncio.get_event_loop()
            result = await loop.run_in_executor(self.executor, _compress_pdf)
            
            if result["success"]:
                logger.info(f"Compression successful: {result['input_bytes']} -> {result['output_bytes']} bytes ({result['ratio']:.1f}% reduction)")
                # Add output path to result
                result["output_path"] = str(output_path)
                return result
            else:
                raise ConversionError(f"Compression failed: {result['message']}")
            
        except Exception as e:
            self._cleanup_work_dir(work_dir)
            raise ConversionError(f"PDF compression failed: {str(e)}")
    
    async def extract_text_to_pdf(self, pdf_path: str, job_id: str = None) -> str:
        """Extract text and create a new text-only PDF"""
        if job_id is None:
            job_id = str(uuid.uuid4())
            
        work_dir = self._get_work_dir(job_id)
        
        try:
            # First extract text
            txt_path = await self.pdf_to_txt(pdf_path, job_id=job_id)
            
            # Read extracted text
            with open(txt_path, "r", encoding="utf-8") as f:
                text_content = f.read()
            
            # Create text-only PDF using reportlab if available
            try:
                from reportlab.pdfgen import canvas
                from reportlab.lib.pagesizes import letter, A4
                from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
                from reportlab.lib.styles import getSampleStyleSheet
                from reportlab.lib.units import inch
                
                output_path = work_dir / "text_only.pdf"
                doc = SimpleDocTemplate(str(output_path), pagesize=A4)
                styles = getSampleStyleSheet()
                story = []
                
                # Split text into paragraphs
                paragraphs = text_content.split('\n\n')
                for para_text in paragraphs:
                    if para_text.strip():
                        para = Paragraph(para_text.strip(), styles['Normal'])
                        story.append(para)
                        story.append(Spacer(1, 12))
                
                doc.build(story)
                return str(output_path)
                
            except ImportError:
                # Fallback: convert text to LibreOffice document then to PDF
                txt_doc_path = work_dir / "extracted_text.odt"
                
                # Create ODT document with LibreOffice
                cmd = [
                    "soffice", "--headless", "--invisible", "--nodefault",
                    "--norestore", "--nologo", "--convert-to", "odt",
                    "--outdir", str(work_dir), txt_path
                ]
                
                subprocess.run(cmd, capture_output=True, timeout=30)
                
                # Convert ODT to PDF
                output_path = await self.docx_to_pdf(str(txt_doc_path), job_id)
                return output_path
            
        except Exception as e:
            self._cleanup_work_dir(work_dir)
            raise ConversionError(f"Text-to-PDF conversion failed: {str(e)}")
    
    async def create_ico_from_pdf(self, pdf_path: str, sizes: List[int] = None, job_id: str = None) -> str:
        """Create ICO file from first page of PDF"""
        if sizes is None:
            sizes = [16, 32, 48, 64, 128, 256]
            
        if job_id is None:
            job_id = str(uuid.uuid4())
            
        work_dir = self._get_work_dir(job_id)
        
        try:
            # First convert PDF to PNG
            png_paths = await self.pdf_to_images(pdf_path, fmt="png", dpi=300, job_id=job_id)
            if not png_paths:
                raise ConversionError("No pages found in PDF")
            
            # Use first page
            first_page_path = png_paths[0]
            
            def _create_ico():
                with Image.open(first_page_path) as img:
                    # Create ICO with multiple sizes
                    ico_images = []
                    for size in sizes:
                        resized = img.resize((size, size), Image.Resampling.LANCZOS)
                        ico_images.append(resized)
                    
                    output_path = work_dir / "icon.ico"
                    ico_images[0].save(
                        str(output_path),
                        format='ICO',
                        sizes=[(size, size) for size in sizes]
                    )
                    
                    return str(output_path)
            
            loop = asyncio.get_event_loop()
            result_path = await loop.run_in_executor(self.executor, _create_ico)
            return result_path
            
        except Exception as e:
            self._cleanup_work_dir(work_dir)
            raise ConversionError(f"ICO creation failed: {str(e)}")


# Global instance
conversion_service = ConversionService()